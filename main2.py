from pptx import Presentation
from pptx.util import Inches


# Создаем объект презентации
prs = Presentation()

# Добавляем слайд

slide_layout = prs.slide_layouts[0]  # Выбираем макет слайда

slide = prs.slides.add_slide(slide_layout)

# Добавляем заголовок и подзаголовок на слайд
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Привет, мир!"
subtitle.text = "Это моя первая презентация с использованием Python"


# Добавляем новый слайд для изображения
slide_layout = prs.slide_layouts[5]  # Пустой слайд
slide = prs.slides.add_slide(slide_layout)

# Добавляем изображение
img_path_1 = '1.jpg'
img_path_2 = '2.jpg'
left = top = Inches(1)
slide.shapes.add_picture(img_path_1, left, top)
slide.shapes.add_picture(img_path_2, left, top)

# Сохраняем презентацию
prs.save('test_presentation_01.pptx')